"""
File upload and processing
"""
import os
import struct
import traceback
import zipfile
from itertools import islice
import pandas as pd
from utils.error_handler import FileError, error_logger
from utils.token_counter import TokenCounter

class FileHandler:
    """Handle various file formats"""
    
    SUPPORTED_FORMATS = {
        "text": [".txt", ".md", ".csv", ".log"],
        "code": [".py", ".js", ".java", ".cpp", ".html", ".css", ".json", ".sql", ".sh"],
        "pdf": [".pdf"],
        "excel": [".xlsx", ".xls"],
        "word": [".docx"],
        "image": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"],
        "powerpoint": [".pptx"]
    }
    
    MAX_FILES = 5
    MAX_SIZE = 10_000_000  # 10MB
    MAX_OOXML_MEMBERS = 4_096
    MAX_OOXML_UNCOMPRESSED_SIZE = 100_000_000
    BINARY_FORMATS = {"pdf", "excel", "word", "image", "powerpoint"}
    OOXML_RESOURCE_LIMIT_EXCEEDED = "ooxml_resource_limit_exceeded"
    RESOURCE_LIMIT_MESSAGE = "File exceeds supported processing limits."

    OOXML_REQUIRED_MEMBERS = {
        ".xlsx": "xl/workbook.xml",
        ".docx": "word/document.xml",
        ".pptx": "ppt/presentation.xml",
    }

    IMAGE_SIGNATURES = {
        ".jpg": lambda header: header.startswith(b"\xff\xd8\xff"),
        ".jpeg": lambda header: header.startswith(b"\xff\xd8\xff"),
        ".png": lambda header: header.startswith(b"\x89PNG\r\n\x1a\n"),
        ".gif": lambda header: header.startswith((b"GIF87a", b"GIF89a")),
        ".bmp": lambda header: header.startswith(b"BM"),
        ".webp": lambda header: header.startswith(b"RIFF") and header[8:12] == b"WEBP",
    }

    CFBF_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
    CFBF_FREE_SECTOR = 0xFFFFFFFF
    CFBF_END_OF_CHAIN = 0xFFFFFFFE
    CFBF_FAT_SECTOR = 0xFFFFFFFD
    MAX_XLS_DIRECTORY_SECTORS = 128
    PARSER_ERROR_MESSAGES = {
        "text": "File could not be read as UTF-8 text.",
        "code": "File could not be read as UTF-8 text.",
        "pdf": "File could not be read. It may be malformed or use unsupported features.",
        "excel": "File could not be read. It may be malformed or use unsupported features.",
        "word": "File could not be read. It may be malformed or use unsupported features.",
        "powerpoint": "File could not be read. It may be malformed or use unsupported features.",
        "image": "Image analysis failed. Please try again.",
    }
    
    @classmethod
    def get_format(cls, filename):
        """Detect file format from extension"""
        ext = os.path.splitext(filename)[1].lower()
        for fmt, extensions in cls.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return fmt
        return None

    @classmethod
    def _is_valid_binary_candidate(cls, file, fmt):
        """Perform lightweight, non-executing validation before binary dispatch."""
        extension = os.path.splitext(file.name)[1].lower()
        try:
            original_position = file.tell()
            file.seek(0)

            if fmt == "pdf":
                header = file.read(1024)
                file.seek(0, os.SEEK_END)
                tail_start = max(file.tell() - 1024, 0)
                file.seek(tail_start)
                return b"%PDF-" in header and b"%%EOF" in file.read(1024)

            if extension == ".xls":
                return cls._has_legacy_excel_workbook_stream(file)

            if extension in cls.OOXML_REQUIRED_MEMBERS:
                return cls._validate_ooxml_archive(file, extension)

            if fmt == "image":
                return cls.IMAGE_SIGNATURES[extension](file.read(12))

            return False
        except (AttributeError, OSError, ValueError, zipfile.BadZipFile):
            return False
        finally:
            try:
                file.seek(original_position)
            except (AttributeError, OSError, ValueError, UnboundLocalError):
                pass

    @classmethod
    def _validate_ooxml_archive(cls, file, extension):
        """Validate OOXML ZIP metadata and required structural members."""
        with zipfile.ZipFile(file) as archive:
            members = archive.infolist()
            if len(members) > cls.MAX_OOXML_MEMBERS:
                return cls.OOXML_RESOURCE_LIMIT_EXCEEDED

            total_uncompressed_size = 0
            member_names = set()
            for member in members:
                total_uncompressed_size += member.file_size
                if total_uncompressed_size > cls.MAX_OOXML_UNCOMPRESSED_SIZE:
                    return cls.OOXML_RESOURCE_LIMIT_EXCEEDED
                member_names.add(member.filename)

        return (
            "[Content_Types].xml" in member_names
            and cls.OOXML_REQUIRED_MEMBERS[extension] in member_names
        )

    @classmethod
    def _has_legacy_excel_workbook_stream(cls, file):
        """Boundedly inspect a CFBF directory for a legacy Excel workbook stream."""
        header = file.read(512)
        if (
            len(header) != 512
            or header[:8] != cls.CFBF_SIGNATURE
            or header[28:30] != b"\xfe\xff"
        ):
            return False

        sector_shift = struct.unpack_from("<H", header, 30)[0]
        mini_sector_shift = struct.unpack_from("<H", header, 32)[0]
        if sector_shift not in (9, 12) or mini_sector_shift != 6:
            return False

        sector_size = 1 << sector_shift
        fat_sector_count = struct.unpack_from("<I", header, 44)[0]
        first_directory_sector = struct.unpack_from("<I", header, 48)[0]
        if not 0 < fat_sector_count <= 109:
            return False

        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        sector_count = (file_size - 512) // sector_size
        if sector_count <= 0 or first_directory_sector >= sector_count:
            return False

        difat_entries = struct.unpack_from("<109I", header, 76)
        fat_sector_ids = difat_entries[:fat_sector_count]
        if any(sector_id >= sector_count for sector_id in fat_sector_ids):
            return False

        fat_entries = []
        for sector_id in fat_sector_ids:
            file.seek(512 + sector_id * sector_size)
            fat_sector = file.read(sector_size)
            if len(fat_sector) != sector_size:
                return False
            fat_entries.extend(struct.unpack(f"<{sector_size // 4}I", fat_sector))

        directory_sector = first_directory_sector
        for _ in range(cls.MAX_XLS_DIRECTORY_SECTORS):
            if directory_sector >= sector_count or directory_sector >= len(fat_entries):
                return False

            file.seek(512 + directory_sector * sector_size)
            directory_data = file.read(sector_size)
            if len(directory_data) != sector_size:
                return False

            for offset in range(0, sector_size, 128):
                name_length = struct.unpack_from("<H", directory_data, offset + 64)[0]
                object_type = directory_data[offset + 66]
                if object_type != 2 or not 2 <= name_length <= 64 or name_length % 2:
                    continue
                name = directory_data[offset:offset + name_length - 2].decode("utf-16le", "ignore")
                if name in {"Workbook", "Book"}:
                    return True

            directory_sector = fat_entries[directory_sector]
            if directory_sector == cls.CFBF_END_OF_CHAIN:
                return False
            if directory_sector in (cls.CFBF_FREE_SECTOR, cls.CFBF_FAT_SECTOR):
                return False

        return False
    
    @classmethod
    def handle(cls, uploaded_files, gemini_agent=None):
        """Process uploaded files"""
        if len(uploaded_files) > cls.MAX_FILES:
            raise FileError(f"Max {cls.MAX_FILES} files allowed")
        
        results = []
        total_tokens = 0
        
        for file in uploaded_files:
            if file.size > cls.MAX_SIZE:
                results.append({
                    "filename": file.name,
                    "error": f"File too large (max 10MB)"
                })
                continue
            
            fmt = cls.get_format(file.name)
            if not fmt:
                results.append({
                    "filename": file.name,
                    "error": "Unsupported format"
                })
                continue
            
            if fmt in cls.BINARY_FORMATS:
                binary_validation = cls._is_valid_binary_candidate(file, fmt)
                if binary_validation == cls.OOXML_RESOURCE_LIMIT_EXCEEDED:
                    results.append({
                        "filename": file.name,
                        "error": cls.RESOURCE_LIMIT_MESSAGE
                    })
                    continue
                if not binary_validation:
                    results.append({
                        "filename": file.name,
                        "error": "Invalid or mismatched binary file"
                    })
                    continue

            try:
                content = cls._process_file(file, fmt, gemini_agent)
                tokens = TokenCounter.count(content) if content else 0
                total_tokens += tokens
                
                results.append({
                    "filename": file.name,
                    "format": fmt,
                    "content": content[:10000],  # Limit
                    "tokens": tokens,
                    "size_kb": file.size / 1024
                })
            
            except Exception as e:
                error_logger.log(
                    "FILE_PROCESSING_ERROR",
                    f"{fmt} processing failed: {type(e).__name__}",
                    details=traceback.format_exc(),
                )
                results.append({
                    "filename": file.name,
                    "error": cls.PARSER_ERROR_MESSAGES.get(
                        fmt,
                        "File could not be processed. Please try again.",
                    )
                })
        
        return {
            "files": results,
            "total_tokens": total_tokens,
            "count": len(results)
        }
    
    @classmethod
    def _process_file(cls, file, fmt, gemini_agent=None):
        """Process file based on format"""
        
        if fmt == "text":
            return file.read().decode('utf-8')[:50000]
        
        elif fmt == "code":
            content = file.read().decode('utf-8')
            ext = os.path.splitext(file.name)[1]
            lang_map = {'.py': 'python', '.js': 'javascript', '.java': 'java',
                       '.cpp': 'cpp', '.html': 'html', '.css': 'css',
                       '.json': 'json', '.sql': 'sql', '.sh': 'bash'}
            lang = lang_map.get(ext, '')
            return f"```{lang}\n{content[:30000]}\n```"
        
        elif fmt == "pdf":
            try:
                import pdfplumber
                with pdfplumber.open(file) as pdf:
                    chunks = []
                    remaining = 40000
                    for page in pdf.pages[:10]:
                        page_text = page.extract_text()
                        if page_text:
                            remaining = cls._append_bounded_text(chunks, page_text, remaining)
                            remaining = cls._append_bounded_text(chunks, "\n", remaining)
                            if remaining == 0:
                                break
                    text = "".join(chunks)
                    return text if text else "No text found in PDF"
            except Exception:
                from PyPDF2 import PdfReader
                file.seek(0)
                reader = PdfReader(file)
                chunks = []
                remaining = 40000
                for page in reader.pages[:10]:
                    page_text = page.extract_text()
                    if page_text:
                        remaining = cls._append_bounded_text(chunks, page_text, remaining)
                        remaining = cls._append_bounded_text(chunks, "\n", remaining)
                        if remaining == 0:
                            break
                text = "".join(chunks)
                return text if text else "No text found in PDF"
        
        elif fmt == "excel":
            df = pd.read_excel(file, nrows=100)
            return df.head(100).to_string()[:30000]
        
        elif fmt == "word":
            from docx import Document
            from docx.text.paragraph import Paragraph
            doc = Document(file)
            chunks = []
            remaining = 30000
            for item in doc.iter_inner_content():
                if not isinstance(item, Paragraph):
                    continue
                paragraph_text = item.text
                if not paragraph_text.strip():
                    continue
                if chunks:
                    remaining = cls._append_bounded_text(chunks, "\n", remaining)
                remaining = cls._append_bounded_text(chunks, paragraph_text, remaining)
                if remaining == 0:
                    break
            return "".join(chunks)
        
        elif fmt == "image":
            if gemini_agent:
                result = gemini_agent.analyze_image(file)
                if not isinstance(result, dict) or result.get("status") != "success":
                    raise ValueError("Image analysis did not return a usable response")
                text = result.get("text")
                if not isinstance(text, str) or not text.strip():
                    raise ValueError("Image analysis did not return text")
                return text[:10000]
            else:
                return "[Image - requires Gemini Vision]"
        
        elif fmt == "powerpoint":
            from pptx import Presentation
            prs = Presentation(file)
            chunks = []
            remaining = 30000
            for slide in islice(prs.slides, 20):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        remaining = cls._append_bounded_text(chunks, shape.text, remaining)
                        remaining = cls._append_bounded_text(chunks, "\n", remaining)
                        if remaining == 0:
                            break
                if remaining == 0:
                    break
            return "".join(chunks)
        
        return None

    @staticmethod
    def _append_bounded_text(chunks, text, remaining):
        """Append no more than the remaining extraction budget."""
        if not text or remaining <= 0:
            return remaining
        chunk = text[:remaining]
        chunks.append(chunk)
        return remaining - len(chunk)
